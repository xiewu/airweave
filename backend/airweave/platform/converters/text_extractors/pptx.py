"""Direct text extraction from PPTX files using python-pptx.

Extracts text from slide shapes, tables, and notes to produce markdown.
Images and diagrams are not captured -- this is a text-only extraction.
"""

from __future__ import annotations

import asyncio
import os
from typing import Any, Optional

from airweave.core.logging import logger
from airweave.platform.sync.exceptions import SyncFailureError

# Minimum total characters to consider the extraction successful.
MIN_TOTAL_CHARS = 50


def _extract_shape_text(shape: Any) -> list[str]:
    """Extract text lines from a single PPTX shape.

    Handles both text-frame shapes and table shapes.

    Args:
        shape: A ``pptx.shapes.base.BaseShape`` instance.

    Returns:
        List of text lines (may be empty).
    """
    lines: list[str] = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)

    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

    return lines


def _extract_slide(slide: Any, slide_idx: int) -> str:
    """Extract markdown for a single slide.

    Args:
        slide: A ``pptx.slide.Slide`` instance.
        slide_idx: 1-based slide number (used for the heading).

    Returns:
        Markdown string for the slide.
    """
    parts: list[str] = [f"## Slide {slide_idx}"]

    for shape in slide.shapes:
        parts.extend(_extract_shape_text(shape))

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            parts.append(f"\n> **Notes:** {notes_text}")

    return "\n\n".join(parts)


async def extract_pptx_text(path: str) -> Optional[str]:
    """Extract text from a PPTX and return markdown.

    Iterates over slides, shapes, tables, and notes to produce a markdown
    representation.  Images and diagrams are not captured.

    Args:
        path: Path to the PPTX file.

    Returns:
        Markdown string if extraction yielded sufficient text, ``None`` otherwise.

    Raises:
        SyncFailureError: If python-pptx is not installed.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise SyncFailureError("python-pptx required for PPTX text extraction but not installed")

    def _extract() -> Optional[str]:
        name = os.path.basename(path)

        try:
            prs = Presentation(path)
        except Exception as exc:
            logger.warning(f"Failed to open PPTX {name}: {exc}")
            return None

        slide_markdowns = [
            _extract_slide(slide, idx) for idx, slide in enumerate(prs.slides, start=1)
        ]
        markdown = "\n\n---\n\n".join(slide_markdowns)

        total_chars = len(markdown.strip())
        if total_chars < MIN_TOTAL_CHARS:
            logger.debug(f"PPTX {name}: only {total_chars} chars extracted, insufficient")
            return None

        logger.debug(f"PPTX {name}: extracted {total_chars} chars")
        return markdown

    return await asyncio.to_thread(_extract)
