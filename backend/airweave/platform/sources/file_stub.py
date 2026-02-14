"""File stub source for testing document conversion pipelines.

Generates exactly one of each supported file type:
- Born-digital PDF (text layer, no OCR needed)
- Scanned PDF (image-only pages, requires OCR)
- PPTX presentation
- DOCX document

Each file embeds an optional ``custom_content_prefix`` (tracking token) so
tests can search for that string to verify end-to-end extraction.
"""

import io
import os
import random
import tempfile
from datetime import datetime, timedelta
from typing import Any, AsyncGenerator, Dict, List, Optional, Tuple

from airweave.platform.configs.auth import FileStubAuthConfig
from airweave.platform.configs.config import FileStubConfig
from airweave.platform.decorators import source
from airweave.platform.entities._base import BaseEntity, Breadcrumb
from airweave.platform.entities.file_stub import (
    DocxFileStubEntity,
    FileStubContainerEntity,
    PdfFileStubEntity,
    PptxFileStubEntity,
    ScannedPdfFileStubEntity,
)
from airweave.platform.sources._base import BaseSource
from airweave.schemas.source_connection import AuthenticationMethod

# Word lists for deterministic content generation (shared with stub.py)
NOUNS = [
    "project",
    "task",
    "document",
    "report",
    "meeting",
    "analysis",
    "review",
    "strategy",
    "plan",
    "update",
    "milestone",
    "feature",
    "module",
    "component",
]
VERBS = [
    "implement",
    "review",
    "update",
    "create",
    "analyze",
    "test",
    "deploy",
    "configure",
    "optimize",
    "refactor",
    "debug",
    "document",
    "design",
]
ADJECTIVES = [
    "important",
    "urgent",
    "critical",
    "minor",
    "major",
    "quick",
    "detailed",
    "comprehensive",
    "preliminary",
    "final",
    "draft",
    "approved",
    "pending",
]
AUTHORS = [
    "Alice Smith",
    "Bob Johnson",
    "Carol Williams",
    "David Brown",
    "Eve Davis",
]


def _load_truetype_font(size: int):
    """Load a truetype font that works across platforms.

    Tries common system font paths for macOS, Linux (Debian/Ubuntu), and
    Windows.  Falls back to Pillow's built-in ``load_default(size=...)``
    (requires Pillow >= 10.1) so CI environments always get a real
    sized font instead of a tiny bitmap.
    """
    from PIL import ImageFont

    # Common TrueType font paths by platform
    _FONT_CANDIDATES = [
        # macOS
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        # Linux (Debian/Ubuntu packages: fonts-dejavu-core, fonts-liberation)
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        # Windows
        "C:\\Windows\\Fonts\\arial.ttf",
    ]

    for path in _FONT_CANDIDATES:
        try:
            return ImageFont.truetype(path, size=size)
        except (IOError, OSError):
            continue

    # Pillow >= 10.1 supports load_default(size=...)
    try:
        return ImageFont.load_default(size=size)
    except TypeError:
        # Pillow < 10.1: load_default() ignores size, but it's the last resort
        return ImageFont.load_default()


class ContentGenerator:
    """Deterministic content generator for file stubs."""

    def __init__(self, seed: int, custom_content_prefix: Optional[str] = None):
        """Initialize the generator with a seed and optional tracking prefix."""
        self.rng = random.Random(seed)
        self.base_time = datetime(2024, 1, 1, 0, 0, 0)
        self.custom_content_prefix = custom_content_prefix

    def _pick(self, items: List[str]) -> str:
        return self.rng.choice(items)

    def _generate_sentence(self, word_count: int = 10) -> str:
        words = []
        for i in range(word_count):
            if i % 3 == 0:
                words.append(self._pick(ADJECTIVES))
            elif i % 3 == 1:
                words.append(self._pick(NOUNS))
            else:
                words.append(self._pick(VERBS))
        return " ".join(words).capitalize() + "."

    def _generate_paragraph(self, sentence_count: int = 5) -> str:
        return " ".join(
            self._generate_sentence(self.rng.randint(8, 15)) for _ in range(sentence_count)
        )

    def _generate_timestamp(self, days_offset: int = 0) -> datetime:
        return self.base_time + timedelta(
            days=days_offset,
            hours=self.rng.randint(0, 23),
            minutes=self.rng.randint(0, 59),
        )

    def generate_title(self) -> str:
        """Generate a random document title."""
        adj = self._pick(ADJECTIVES)
        noun = self._pick(NOUNS)
        verb = self._pick(VERBS)
        return f"{adj.capitalize()} {noun} to {verb}"

    # ── Born-digital PDF ─────────────────────────────────────────────────
    def generate_pdf_content(self) -> Tuple[bytes, int]:
        """Generate a born-digital PDF with embedded text.

        Returns:
            Tuple of (pdf_bytes, page_count).
        """
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        pdf.add_page()
        pdf.set_font("Helvetica", "B", 22)
        title = self.generate_title()
        pdf.cell(0, 14, title, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(6)

        if self.custom_content_prefix:
            pdf.set_font("Helvetica", size=14)
            pdf.multi_cell(0, 9, self.custom_content_prefix)
            pdf.ln(6)

        num_sections = self.rng.randint(4, 7)
        for i in range(num_sections):
            section_title = f"Section {i + 1}: {self.generate_title()}"
            pdf.set_font("Helvetica", "B", 16)
            pdf.cell(0, 12, section_title, new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", size=13)
            num_paragraphs = self.rng.randint(2, 3)
            for _ in range(num_paragraphs):
                paragraph = self._generate_paragraph(self.rng.randint(6, 10))
                pdf.multi_cell(0, 8, paragraph)
                pdf.ln(4)
            pdf.ln(5)

        page_count = pdf.page
        return bytes(pdf.output()), page_count

    # ── Scanned (image-only) PDF ─────────────────────────────────────────
    def generate_scanned_pdf_content(self) -> Tuple[bytes, int]:
        """Generate an image-only PDF (no text layer).

        Renders text onto images using Pillow, then embeds those images
        into a PDF with fpdf2.  The resulting PDF has **no** text layer
        and requires OCR to extract content.

        Returns:
            Tuple of (pdf_bytes, page_count).
        """
        from fpdf import FPDF
        from PIL import Image, ImageDraw

        pdf = FPDF()
        pdf.set_auto_page_break(auto=False)

        # Page dimensions in points (A4-ish)
        page_w_pt, page_h_pt = 595, 842
        dpi = 300
        img_w = int(page_w_pt * dpi / 72)
        img_h = int(page_h_pt * dpi / 72)

        num_pages = self.rng.randint(1, 3)

        for page_idx in range(num_pages):
            # Create a white image
            img = Image.new("RGB", (img_w, img_h), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # At 300 DPI: 67px ≈ 16pt, 54px ≈ 13pt
            # Try common system fonts across platforms (macOS, Linux, Windows).
            # Pillow's load_default() returns a tiny bitmap font that ignores
            # size, so we MUST find a real truetype font for OCR-readable text.
            font = _load_truetype_font(size=67)
            small_font = _load_truetype_font(size=54)

            y = 120

            # Title
            if page_idx == 0:
                title = self.generate_title()
                draw.text((120, y), title, fill=(0, 0, 0), font=font)
                y += 100

                # Embed the tracking token prominently
                if self.custom_content_prefix:
                    draw.text((120, y), self.custom_content_prefix, fill=(0, 0, 0), font=font)
                    y += 100

            # Section heading
            section_title = f"Section {page_idx + 1}: {self.generate_title()}"
            draw.text((120, y), section_title, fill=(0, 0, 0), font=font)
            y += 90

            # Body text lines
            for _ in range(self.rng.randint(8, 15)):
                line = self._generate_sentence(self.rng.randint(6, 12))
                # Truncate to fit page width
                draw.text((120, y), line[:80], fill=(30, 30, 30), font=small_font)
                y += 70
                if y > img_h - 160:
                    break

            # Save image to temp file for fpdf
            img_buf = io.BytesIO()
            img.save(img_buf, format="JPEG", quality=95)
            img_buf.seek(0)

            # Write temp jpeg file (fpdf needs a path)
            tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
            tmp.write(img_buf.getvalue())
            tmp.close()

            pdf.add_page()
            pdf.image(tmp.name, x=0, y=0, w=page_w_pt * 25.4 / 72, h=page_h_pt * 25.4 / 72)
            os.unlink(tmp.name)

        page_count = pdf.page
        return bytes(pdf.output()), page_count

    # ── PPTX ─────────────────────────────────────────────────────────────
    def generate_pptx_content(self) -> Tuple[bytes, int]:
        """Generate a PPTX with text content.

        Returns:
            Tuple of (pptx_bytes, slide_count).
        """
        from pptx import Presentation

        prs = Presentation()

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = self.generate_title()
        subtitle = slide.placeholders[1]
        if self.custom_content_prefix:
            subtitle.text = self.custom_content_prefix
        else:
            subtitle.text = self._generate_sentence(8)

        # Content slides
        num_slides = self.rng.randint(2, 4)
        for i in range(num_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Section {i + 1}: {self.generate_title()}"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = self._generate_paragraph(self.rng.randint(3, 5))
            p = tf.add_paragraph()
            p.text = self._generate_paragraph(self.rng.randint(2, 4))

        slide_count = len(prs.slides)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), slide_count

    # ── DOCX ─────────────────────────────────────────────────────────────
    def generate_docx_content(self) -> Tuple[bytes, int]:
        """Generate a DOCX with text content.

        Returns:
            Tuple of (docx_bytes, approximate_page_count).
        """
        from docx import Document

        doc = Document()

        doc.add_heading(self.generate_title(), level=1)

        if self.custom_content_prefix:
            doc.add_paragraph(self.custom_content_prefix)

        num_sections = self.rng.randint(2, 4)
        for i in range(num_sections):
            doc.add_heading(f"Section {i + 1}: {self.generate_title()}", level=2)
            for _ in range(self.rng.randint(2, 4)):
                doc.add_paragraph(self._generate_paragraph(self.rng.randint(3, 6)))

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue(), num_sections  # rough page estimate


@source(
    name="File Stub",
    short_name="file_stub",
    auth_methods=[AuthenticationMethod.DIRECT],
    oauth_type=None,
    auth_config_class=FileStubAuthConfig,
    config_class=FileStubConfig,
    labels=["Internal", "Testing"],
    supports_continuous=False,
    internal=True,
)
class FileStubSource(BaseSource):
    """File stub source for testing document conversion pipelines.

    Generates exactly one of each file type: born-digital PDF, scanned PDF,
    PPTX, and DOCX. All files embed the tracking token for search assertions.
    """

    def __init__(self):
        """Initialize with default seed and no tracking prefix."""
        super().__init__()
        self.seed: int = 42
        self.custom_content_prefix: Optional[str] = None
        self.generator: Optional[ContentGenerator] = None
        self._temp_dir: Optional[str] = None

    @classmethod
    async def create(
        cls,
        credentials: Optional[FileStubAuthConfig] = None,
        config: Optional[Dict[str, Any]] = None,
    ) -> "FileStubSource":
        """Create a configured FileStubSource from credentials and config."""
        instance = cls()
        config = config or {}

        instance.seed = config.get("seed", 42)
        instance.custom_content_prefix = config.get("custom_content_prefix", None)
        instance.generator = ContentGenerator(
            seed=instance.seed,
            custom_content_prefix=instance.custom_content_prefix,
        )
        return instance

    async def _write_binary(self, data: bytes, extension: str, name: str) -> str:
        """Write binary data to a temp file and return its path."""
        if self._temp_dir is None:
            self._temp_dir = tempfile.mkdtemp(prefix="file_stub_")
        filepath = os.path.join(self._temp_dir, f"{name}{extension}")
        with open(filepath, "wb") as f:
            f.write(data)
        return filepath

    async def generate_entities(self) -> AsyncGenerator[BaseEntity, None]:
        """Yield a container + one entity per file type."""
        gen = self.generator
        assert gen is not None

        self.logger.info(f"FileStubSource: seed={self.seed}")

        # Container
        container_id = f"file-stub-container-{self.seed}"
        container = FileStubContainerEntity(
            container_id=container_id,
            container_name=f"File Stub Container (seed={self.seed})",
            description="Container for file converter pipeline tests",
            created_at=datetime(2024, 1, 1, 0, 0, 0),
            seed=self.seed,
            entity_count=4,
            breadcrumbs=[],
        )
        yield container

        breadcrumbs = [
            Breadcrumb(
                entity_id=container_id,
                name=container.container_name,
                entity_type="FileStubContainerEntity",
            )
        ]

        idx = 0

        # ── 1. Born-digital PDF ──────────────────────────────────────────
        pdf_bytes, page_count = gen.generate_pdf_content()
        name = f"born_digital_pdf_{self.seed}"
        filepath = await self._write_binary(pdf_bytes, ".pdf", name)
        yield PdfFileStubEntity(
            stub_id=f"born-digital-pdf-{self.seed}",
            file_name=f"{name}.pdf",
            description="Born-digital PDF with embedded text layer",
            author=gen._pick(AUTHORS),
            page_count=page_count,
            created_at=gen._generate_timestamp(idx),
            modified_at=gen._generate_timestamp(idx + 1),
            sequence_number=idx,
            breadcrumbs=breadcrumbs,
            url=f"file-stub://pdf/{name}.pdf",
            size=len(pdf_bytes),
            file_type="document",
            mime_type="application/pdf",
            local_path=filepath,
        )
        idx += 1

        # ── 2. Scanned (image-only) PDF ─────────────────────────────────
        scan_bytes, scan_pages = gen.generate_scanned_pdf_content()
        name = f"scanned_pdf_{self.seed}"
        filepath = await self._write_binary(scan_bytes, ".pdf", name)
        yield ScannedPdfFileStubEntity(
            stub_id=f"scanned-pdf-{self.seed}",
            file_name=f"{name}.pdf",
            description="Image-only scanned PDF requiring OCR",
            author=gen._pick(AUTHORS),
            page_count=scan_pages,
            created_at=gen._generate_timestamp(idx),
            modified_at=gen._generate_timestamp(idx + 1),
            sequence_number=idx,
            breadcrumbs=breadcrumbs,
            url=f"file-stub://scanned-pdf/{name}.pdf",
            size=len(scan_bytes),
            file_type="document",
            mime_type="application/pdf",
            local_path=filepath,
        )
        idx += 1

        # ── 3. PPTX ─────────────────────────────────────────────────────
        pptx_bytes, slide_count = gen.generate_pptx_content()
        name = f"presentation_{self.seed}"
        filepath = await self._write_binary(pptx_bytes, ".pptx", name)
        yield PptxFileStubEntity(
            stub_id=f"pptx-{self.seed}",
            file_name=f"{name}.pptx",
            description="PPTX presentation with slide text",
            author=gen._pick(AUTHORS),
            slide_count=slide_count,
            created_at=gen._generate_timestamp(idx),
            modified_at=gen._generate_timestamp(idx + 1),
            sequence_number=idx,
            breadcrumbs=breadcrumbs,
            url=f"file-stub://pptx/{name}.pptx",
            size=len(pptx_bytes),
            file_type="presentation",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            local_path=filepath,
        )
        idx += 1

        # ── 4. DOCX ─────────────────────────────────────────────────────
        docx_bytes, approx_pages = gen.generate_docx_content()
        name = f"document_{self.seed}"
        filepath = await self._write_binary(docx_bytes, ".docx", name)
        yield DocxFileStubEntity(
            stub_id=f"docx-{self.seed}",
            file_name=f"{name}.docx",
            description="DOCX document with paragraph text",
            author=gen._pick(AUTHORS),
            page_count=approx_pages,
            created_at=gen._generate_timestamp(idx),
            modified_at=gen._generate_timestamp(idx + 1),
            sequence_number=idx,
            breadcrumbs=breadcrumbs,
            url=f"file-stub://docx/{name}.docx",
            size=len(docx_bytes),
            file_type="document",
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            local_path=filepath,
        )

        self.logger.info("FileStubSource: generated all 4 file entities")

    async def validate(self) -> bool:
        """Always valid - no external dependencies."""
        return True
