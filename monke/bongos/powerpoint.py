"""PowerPoint bongo implementation.

Creates, updates, and deletes test PowerPoint presentations via the Microsoft Graph API.
"""

import asyncio
import io
import time
import uuid
from typing import Any, Dict, List

import httpx
from monke.bongos.base_bongo import BaseBongo
from monke.generation.powerpoint import generate_presentations
from monke.utils.logging import get_logger

try:
    from pptx import Presentation

    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

GRAPH = "https://graph.microsoft.com/v1.0"


class PowerPointBongo(BaseBongo):
    """Bongo for PowerPoint that creates test entities for E2E testing.

    Key responsibilities:
    - Create test PowerPoint presentations in OneDrive
    - Add content with verification tokens
    - Update presentations to test incremental sync
    - Delete presentations to test deletion detection
    - Clean up all test data
    """

    connector_type = "powerpoint"

    def __init__(self, credentials: Dict[str, Any], **kwargs):
        super().__init__(credentials)
        self.access_token: str = credentials["access_token"]
        self.entity_count: int = int(kwargs.get("entity_count", 3))
        self.openai_model: str = kwargs.get("openai_model", "gpt-4.1-mini")
        self.rate_limit_delay = float(kwargs.get("rate_limit_delay_ms", 500)) / 1000.0
        self.logger = get_logger("powerpoint_bongo")

        self._test_presentations: List[Dict[str, Any]] = []
        self._last_req = 0.0

        if not HAS_PYTHON_PPTX:
            raise ImportError(
                "python-pptx is required for PowerPoint bongo. "
                "Install with: pip install python-pptx"
            )

    async def create_entities(self) -> List[Dict[str, Any]]:
        """Create test PowerPoint presentations in OneDrive."""
        self.logger.info(
            f"🥁 Creating {self.entity_count} PowerPoint test presentations"
        )
        out: List[Dict[str, Any]] = []

        tokens = [uuid.uuid4().hex[:8] for _ in range(self.entity_count)]
        test_name = f"Monke_PPT_{uuid.uuid4().hex[:8]}"
        presentations = await generate_presentations(
            self.openai_model, tokens, test_name
        )

        self.logger.info(f"📝 Generated {len(presentations)} presentations")

        async with httpx.AsyncClient(base_url=GRAPH, timeout=60) as client:
            for pres_data, token in zip(presentations, tokens):
                await self._pace()

                safe_filename = self._sanitize_filename(
                    f"{pres_data.title}.pptx"
                )

                pptx_bytes = self._create_pptx_file(pres_data)

                self.logger.info(f"📤 Uploading PowerPoint: {safe_filename}")
                upload_url = f"/me/drive/root:/{safe_filename}:/content"

                r = await client.put(
                    upload_url,
                    headers={
                        "Authorization": f"Bearer {self.access_token}",
                        "Content-Type": (
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation"
                        ),
                    },
                    content=pptx_bytes,
                )

                if r.status_code not in (200, 201):
                    self.logger.error(f"Upload failed {r.status_code}: {r.text}")
                    r.raise_for_status()

                doc_file = r.json()
                doc_id = doc_file["id"]

                self.logger.info(f"✅ Uploaded presentation: {doc_id} - {safe_filename}")

                ent = {
                    "type": "presentation",
                    "id": doc_id,
                    "filename": safe_filename,
                    "title": pres_data.title,
                    "token": token,
                    "expected_content": token,
                }
                out.append(ent)
                self._test_presentations.append(ent)
                self.created_entities.append({"id": doc_id, "name": safe_filename})

                self.logger.info(
                    f"📝 Presentation '{safe_filename}' created with token: {token}"
                )

        self.logger.info(
            f"✅ Created {len(self._test_presentations)} PowerPoint presentations"
        )
        return out

    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename for OneDrive (remove illegal characters)."""
        illegal_chars = ['\\', '/', ':', '*', '?', '"', '<', '>', '|']
        safe_name = filename
        for char in illegal_chars:
            safe_name = safe_name.replace(char, "_")
        safe_name = safe_name.strip(". ")
        if len(safe_name) > 200:
            name, ext = safe_name.rsplit(".", 1) if "." in safe_name else (safe_name, "")
            safe_name = name[:195] + "." + ext if ext else name[:200]
        return safe_name

    def _create_pptx_file(self, pres_data: Any) -> bytes:
        """Create a PowerPoint file with the given content.

        Args:
            pres_data: PowerPointPresentationContent with title and content

        Returns:
            Bytes of the .pptx file
        """
        prs = Presentation()

        # Split content by slide separator
        slide_contents = pres_data.content.split("---")
        slide_contents = [s.strip() for s in slide_contents if s.strip()]
        if not slide_contents:
            slide_contents = [pres_data.content or "Slide content"]

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = pres_data.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_contents[0][:500]

        # Content slides
        for i, content in enumerate(slide_contents[1:], start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i + 1}"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = content[:2000]

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    async def update_entities(self) -> List[Dict[str, Any]]:
        """Update PowerPoint presentations by adding a slide with the same token."""
        if not self._test_presentations:
            return []

        self.logger.info(
            f"🥁 Updating {min(2, len(self._test_presentations))} "
            "PowerPoint presentations"
        )
        updated = []

        async with httpx.AsyncClient(base_url=GRAPH, timeout=60) as client:
            for ent in self._test_presentations[: min(2, len(self._test_presentations))]:
                await self._pace()

                download_url = f"/me/drive/items/{ent['id']}/content"
                r = await client.get(download_url, headers=self._hdrs())

                if r.status_code != 200:
                    self.logger.warning(
                        f"Failed to download presentation: {r.status_code} - {r.text[:200]}"
                    )
                    continue

                prs = Presentation(io.BytesIO(r.content))
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Update Section"
                body = slide.placeholders[1]
                body.text_frame.text = (
                    f"This presentation was updated. Token: {ent['token']}"
                )

                buf = io.BytesIO()
                prs.save(buf)
                updated_bytes = buf.getvalue()

                upload_url = f"/me/drive/items/{ent['id']}/content"
                r = await client.put(
                    upload_url,
                    headers={
                        "Authorization": f"Bearer {self.access_token}",
                        "Content-Type": (
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation"
                        ),
                    },
                    content=updated_bytes,
                )

                if r.status_code in (200, 201):
                    updated.append({**ent, "updated": True})
                    self.logger.info(
                        f"📝 Updated presentation '{ent['filename']}' "
                        f"with token: {ent['token']}"
                    )
                else:
                    self.logger.warning(
                        f"Failed to update: {r.status_code} - {r.text[:200]}"
                    )

                await asyncio.sleep(0.5)

        return updated

    async def delete_entities(self) -> List[str]:
        """Delete all test presentations."""
        return await self.delete_specific_entities(self._test_presentations.copy())

    async def delete_specific_entities(
        self, entities: List[Dict[str, Any]]
    ) -> List[str]:
        """Delete specific PowerPoint presentations with retry for locked files."""
        if not entities:
            entities = list(self._test_presentations)
        if not entities:
            return []

        self.logger.info(f"🥁 Deleting {len(entities)} PowerPoint presentations")
        deleted: List[str] = []

        # Iterate over a copy so we can safely mutate self._test_presentations in the loop
        to_delete = list(entities)
        async with httpx.AsyncClient(base_url=GRAPH, timeout=30) as client:
            for ent in to_delete:
                try:
                    await self._pace()
                    max_retries = 3
                    retry_delay = 2.0

                    for attempt in range(max_retries):
                        r = await client.delete(
                            f"/me/drive/items/{ent['id']}",
                            headers=self._hdrs(),
                        )

                        if r.status_code == 204:
                            deleted.append(ent["id"])
                            self.logger.info(
                                f"✅ Deleted: {ent.get('filename', ent['id'])}"
                            )
                            if ent in self._test_presentations:
                                self._test_presentations.remove(ent)
                            break

                        if r.status_code == 423 and attempt < max_retries - 1:
                            self.logger.warning(
                                f"⏳ Presentation locked (423), retrying in "
                                f"{retry_delay}s (attempt {attempt + 1}/{max_retries})"
                            )
                            await asyncio.sleep(retry_delay)
                            retry_delay *= 2
                        else:
                            self.logger.warning(
                                f"Delete failed: {r.status_code} - {r.text[:200]}"
                            )
                            break

                except Exception as e:
                    self.logger.warning(
                        f"Delete error for {ent.get('filename', ent['id'])}: {e}"
                    )

        return deleted

    async def cleanup(self):
        """Comprehensive cleanup of all test resources."""
        self.logger.info("🧹 Starting comprehensive PowerPoint cleanup")

        cleanup_stats = {"presentations_deleted": 0, "errors": 0}

        try:
            async with httpx.AsyncClient(base_url=GRAPH, timeout=30) as client:
                if self._test_presentations:
                    self.logger.info(
                        f"🗑️ Deleting {len(self._test_presentations)} "
                        "tracked presentations"
                    )
                    deleted = await self.delete_specific_entities(
                        self._test_presentations.copy()
                    )
                    cleanup_stats["presentations_deleted"] += len(deleted)

                await self._cleanup_orphaned_presentations(client, cleanup_stats)

            self.logger.info(
                f"🧹 Cleanup completed: {cleanup_stats['presentations_deleted']} "
                f"presentations deleted, {cleanup_stats['errors']} errors"
            )
        except Exception as e:
            self.logger.error(f"❌ Error during comprehensive cleanup: {e}")

    async def _cleanup_orphaned_presentations(
        self, client: httpx.AsyncClient, stats: Dict[str, Any]
    ):
        """Find and delete orphaned test presentations from previous runs."""
        try:
            await self._pace()
            r = await client.get(
                "/me/drive/root/children",
                headers=self._hdrs(),
            )

            if r.status_code == 200:
                files = r.json().get("value", [])
                test_files = [
                    f
                    for f in files
                    if f.get("name", "").startswith("Monke_PPT_")
                    and f.get("name", "").endswith(".pptx")
                ]

                if test_files:
                    self.logger.info(
                        f"🔍 Found {len(test_files)} orphaned test presentations"
                    )
                    for doc in test_files:
                        try:
                            await self._pace()
                            del_r = await client.delete(
                                f"/me/drive/items/{doc['id']}",
                                headers=self._hdrs(),
                            )
                            if del_r.status_code == 204:
                                stats["presentations_deleted"] += 1
                                self.logger.info(
                                    f"✅ Deleted orphaned: {doc.get('name', 'Unknown')}"
                                )
                            else:
                                stats["errors"] += 1
                        except Exception as e:
                            stats["errors"] += 1
                            self.logger.warning(
                                f"⚠️ Failed to delete {doc['id']}: {e}"
                            )
        except Exception as e:
            self.logger.warning(
                f"⚠️ Could not search for orphaned presentations: {e}"
            )

    def _hdrs(self) -> Dict[str, str]:
        """Standard headers for Graph API requests."""
        return {
            "Authorization": f"Bearer {self.access_token}",
            "Content-Type": "application/json",
        }

    async def _pace(self):
        """Rate limiting helper."""
        now = time.time()
        if (delta := now - self._last_req) < self.rate_limit_delay:
            await asyncio.sleep(self.rate_limit_delay - delta)
        self._last_req = time.time()
